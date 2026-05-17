from typing import BinaryIO
import pptx
import pptx.enum.shapes
from .base import BaseConverter, ConversionResult


def _shape_sort_key(shape):
    """Return (top, left) as floats for stable reading-order sort.

    python-pptx position attributes are Emu (int subclass) or None,
    but some exotic shape types return unexpected types. Always coerce
    to float so comparisons never raise TypeError.
    """
    try:
        top = float("-inf") if shape.top is None else float(shape.top)
    except Exception:
        top = float("-inf")
    try:
        left = float("-inf") if shape.left is None else float(shape.left)
    except Exception:
        left = float("-inf")
    return (top, left)


class PptxConverter(BaseConverter):
    def accepts(self, extension: str) -> bool:
        return extension.lower() == ".pptx"

    def convert(self, stream: BinaryIO) -> ConversionResult:
        presentation = pptx.Presentation(stream)
        slides_md = []
        assets: dict[str, bytes] = {}
        img_n: list[int] = [0]

        for i, slide in enumerate(presentation.slides, start=1):
            parts = [f"<!-- Slide {i} -->"]

            shapes = sorted(slide.shapes, key=_shape_sort_key)

            try:
                title_el = slide.shapes.title._element if slide.shapes.title is not None else None
            except Exception:
                title_el = None

            for shape in shapes:
                try:
                    parts.extend(_render_shape(shape, title_el, assets, img_n))
                except Exception:
                    pass

            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame
                if notes and notes.text.strip():
                    parts.append(f"**Notes:** {notes.text.strip()}")

            slides_md.append("\n\n".join(p for p in parts if p))

        return ConversionResult(
            markdown="\n\n---\n\n".join(slides_md),
            assets=assets,
        )


def _is_table(shape) -> bool:
    try:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.TABLE:
            return True
    except Exception:
        pass
    try:
        return shape.has_table
    except Exception:
        return False


def _is_group(shape) -> bool:
    try:
        return shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _render_shape(shape, title_el, assets: dict, img_n: list) -> list[str]:
    parts = []

    if _is_table(shape):
        try:
            md = _table_to_markdown(shape.table)
            if md:
                parts.append(md)
        except Exception:
            pass
        return parts

    # Pictures — extract image bytes and emit a markdown image reference
    try:
        if shape.shape_type == pptx.enum.shapes.MSO_SHAPE_TYPE.PICTURE:
            img_n[0] += 1
            ext = shape.image.ext
            safe_ext = ext if ext and ext.isalnum() and len(ext) <= 10 else "bin"
            filename = f"images/img{img_n[0]:03d}.{safe_ext}"
            assets[filename] = shape.image.blob
            parts.append(f"![]({filename})")
            return parts
    except Exception:
        pass

    try:
        if shape.has_chart:
            chart = shape.chart
            title = ""
            try:
                if chart.has_title:
                    title = f": {chart.chart_title.text_frame.text}"
            except Exception:
                pass

            chart_parts = [f"### Chart{title}"]

            img_bytes = _chart_to_png(chart)
            if img_bytes:
                img_n[0] += 1
                filename = f"images/img{img_n[0]:03d}.png"
                assets[filename] = img_bytes
                chart_parts.append(f"![]({filename})")

            data_table = _chart_data_table(chart)
            if data_table:
                chart_parts.append(data_table)

            parts.append("\n\n".join(chart_parts))
            return parts
    except Exception:
        pass

    if _is_group(shape):
        try:
            sub_shapes = sorted(shape.shapes, key=_shape_sort_key)
            for sub in sub_shapes:
                try:
                    parts.extend(_render_shape(sub, title_el, assets, img_n))
                except Exception:
                    pass
        except Exception:
            pass
        return parts

    try:
        if shape.has_text_frame:
            if title_el is not None and shape._element is title_el:
                title_text = shape.text_frame.text.strip()
                if title_text:
                    parts.append(f"# {title_text}")
            else:
                text = _text_frame_to_md(shape.text_frame)
                if text:
                    parts.append(text)
    except Exception:
        pass

    return parts


_A_NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"


def _para_is_bullet(para) -> bool:
    try:
        pPr = para._p.find(f"{_A_NS}pPr")
        if pPr is None:
            return False
        return (
            pPr.find(f"{_A_NS}buChar") is not None
            or pPr.find(f"{_A_NS}buAutoNum") is not None
        )
    except Exception:
        return False


def _text_frame_to_md(tf) -> str:
    """Convert a text frame to markdown, preserving paragraph structure and bullets.

    Consecutive bullet paragraphs are grouped into a single list block;
    non-bullet paragraphs are separated by blank lines so markdown renders them
    as distinct elements.
    """
    blocks = []
    current_lines: list[str] = []
    current_is_bullet: bool | None = None

    for para in tf.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        level = para.level or 0
        is_bullet = _para_is_bullet(para) or level > 0
        line = ("  " * level + "- " + text) if is_bullet else text

        if current_is_bullet is None or is_bullet == current_is_bullet:
            current_lines.append(line)
            current_is_bullet = is_bullet
        else:
            blocks.append("\n".join(current_lines))
            current_lines = [line]
            current_is_bullet = is_bullet

    if current_lines:
        blocks.append("\n".join(current_lines))

    return "\n\n".join(blocks)


def _table_to_markdown(table) -> str:
    rows = list(table.rows)
    if not rows:
        return ""

    def cell_text(cell):
        return cell.text.replace("|", "\\|").replace("\n", " ").strip()

    header = [cell_text(c) for c in rows[0].cells]
    lines = ["| " + " | ".join(header) + " |",
             "| " + " | ".join("---" for _ in header) + " |"]
    for row in rows[1:]:
        lines.append("| " + " | ".join(cell_text(c) for c in row.cells) + " |")

    return "\n".join(lines)


def _chart_data_table(chart) -> str:
    """Return a markdown table of chart categories × series values."""
    try:
        categories = [c.label for c in chart.plots[0].categories]
        series_names = [s.name for s in chart.series]
        header = ["Category"] + series_names
        rows = []
        for idx, cat in enumerate(categories):
            row = [str(cat)] + [str(s.values[idx]) for s in chart.series]
            rows.append(row)
        lines = ["| " + " | ".join(header) + " |",
                 "| " + " | ".join("---" for _ in header) + " |"]
        for row in rows:
            lines.append("| " + " | ".join(row) + " |")
        return "\n".join(lines)
    except Exception:
        return ""


def _chart_to_png(chart) -> bytes | None:
    """Render chart data to a PNG using matplotlib. Returns bytes or None on failure.

    Uses Figure/FigureCanvasAgg directly (no pyplot global state) so it is safe
    to call from any thread or Streamlit context.
    """
    try:
        import io
        import numpy as np
        from matplotlib.figure import Figure
        from matplotlib.backends.backend_agg import FigureCanvasAgg
        from pptx.enum.chart import XL_CHART_TYPE

        categories = [str(c.label) for c in chart.plots[0].categories]
        series_list = list(chart.series)
        if not series_list or not categories:
            return None

        chart_type = chart.chart_type
        n = len(series_list)
        x = np.arange(len(categories))

        fig = Figure(figsize=(8, 4.5))
        FigureCanvasAgg(fig)
        ax = fig.add_subplot(111)

        _BAR_H = {XL_CHART_TYPE.BAR_CLUSTERED, XL_CHART_TYPE.BAR_STACKED,
                  XL_CHART_TYPE.BAR_STACKED_100}
        _LINE = {XL_CHART_TYPE.LINE, XL_CHART_TYPE.LINE_MARKERS,
                 XL_CHART_TYPE.LINE_STACKED, XL_CHART_TYPE.LINE_MARKERS_STACKED}
        _PIE = {XL_CHART_TYPE.PIE, XL_CHART_TYPE.PIE_EXPLODED}

        def _vals(series):
            return [v if v is not None else 0 for v in series.values]

        def _autoscale_y(ax, all_vals, horizontal=False):
            """Adjust axis bottom so bars look distinct when min >> 0."""
            finite = [v for v in all_vals if v is not None]
            if not finite:
                return
            lo, hi = min(finite), max(finite)
            if lo > 0 and hi > 0 and (hi - lo) < 0.3 * hi:
                pad = (hi - lo) * 0.3 or hi * 0.05
                setter = ax.set_xlim if horizontal else ax.set_ylim
                setter(max(0, lo - pad * 2), hi + pad)

        all_vals = [v for s in series_list for v in _vals(s)]

        if chart_type in _BAR_H:
            width = 0.7 / n
            for i, s in enumerate(series_list):
                offset = (i - n / 2 + 0.5) * width
                ax.barh(x + offset, _vals(s), width, label=s.name or f"S{i+1}")
            ax.set_yticks(x)
            ax.set_yticklabels(categories)
            _autoscale_y(ax, all_vals, horizontal=True)

        elif chart_type in _LINE:
            for i, s in enumerate(series_list):
                ax.plot(x, _vals(s), marker="o", label=s.name or f"S{i+1}")
            ax.set_xticks(x)
            ax.set_xticklabels(categories)

        elif chart_type in _PIE:
            vals = _vals(series_list[0])
            ax.pie(vals, labels=categories, autopct="%1.1f%%")

        else:
            # Default: clustered column (vertical bar)
            width = 0.7 / n
            for i, s in enumerate(series_list):
                offset = (i - n / 2 + 0.5) * width
                ax.bar(x + offset, _vals(s), width, label=s.name or f"S{i+1}")
            ax.set_xticks(x)
            ax.set_xticklabels(categories)
            _autoscale_y(ax, all_vals)

        try:
            if chart.has_title:
                ax.set_title(chart.chart_title.text_frame.text)
        except Exception:
            pass

        if n > 1 or any(s.name for s in series_list):
            ax.legend()

        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=120, bbox_inches="tight")
        buf.seek(0)
        return buf.read()

    except Exception:
        return None
