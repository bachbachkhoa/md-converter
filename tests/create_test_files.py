"""
Run this script once to generate the test files used by the test suite.
    python tests/create_test_files.py
"""
import sys
import os

sys.path.insert(0, os.path.dirname(os.path.dirname(__file__)))

OUT = os.path.join(os.path.dirname(__file__), "test_files")

# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------
def make_docx():
    from docx import Document
    from docx.shared import Pt

    doc = Document()
    doc.add_heading("DOCX-HEADING-a1b2c3d4", level=1)
    doc.add_paragraph("DOCX-PARA-e5f6g7h8")

    doc.add_heading("Section Two", level=2)
    p = doc.add_paragraph()
    p.add_run("DOCX-BOLD-i9j0k1l2").bold = True
    p.add_run(" normal ")
    p.add_run("DOCX-ITALIC-m3n4o5p6").italic = True

    table = doc.add_table(rows=2, cols=3)
    table.rows[0].cells[0].text = "DOCX-TH-COL1"
    table.rows[0].cells[1].text = "DOCX-TH-COL2"
    table.rows[0].cells[2].text = "DOCX-TH-COL3"
    table.rows[1].cells[0].text = "DOCX-TD-R1C1"
    table.rows[1].cells[1].text = "DOCX-TD-R1C2"
    table.rows[1].cells[2].text = "DOCX-TD-R1C3"

    doc.save(os.path.join(OUT, "test.docx"))
    print("Created test.docx")


# ---------------------------------------------------------------------------
# XLSX (with a merged cell)
# ---------------------------------------------------------------------------
def make_xlsx():
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "XLSX-SHEET-q7r8s9t0"

    ws["A1"] = "XLSX-MERGED-HEADER"
    ws.merge_cells("A1:C1")   # merged across 3 columns

    ws["A2"] = "XLSX-TH-COL1"
    ws["B2"] = "XLSX-TH-COL2"
    ws["C2"] = "XLSX-TH-COL3"
    ws["A3"] = "XLSX-TD-u1v2w3"
    ws["B3"] = "XLSX-TD-x4y5z6"
    ws["C3"] = "XLSX-TD-a7b8c9"

    wb.save(os.path.join(OUT, "test.xlsx"))
    print("Created test.xlsx")


# ---------------------------------------------------------------------------
# XLSX (with an embedded image)
# ---------------------------------------------------------------------------
def make_xlsx_with_images():
    import struct
    import zlib
    import openpyxl
    from io import BytesIO
    from openpyxl.drawing.image import Image as XlImage

    def _minimal_png() -> bytes:
        """Build a valid 1×1 white-pixel PNG from raw bytes — no extra deps."""
        def chunk(tag: bytes, data: bytes) -> bytes:
            crc = zlib.crc32(tag + data) & 0xFFFFFFFF
            return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)

        ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
        idat = chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
        iend = chunk(b"IEND", b"")
        return b"\x89PNG\r\n\x1a\n" + ihdr + idat + iend

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "XLSX-IMG-SHEET-r1s2t3u4"

    ws["A1"] = "XLSX-IMG-ROW-BEFORE"
    ws["A2"] = "XLSX-IMG-ROW-AFTER"

    img = XlImage(BytesIO(_minimal_png()))
    ws.add_image(img, "A1")

    wb.save(os.path.join(OUT, "test_xlsx_with_images.xlsx"))
    print("Created test_xlsx_with_images.xlsx")


# ---------------------------------------------------------------------------
# XLSX (with a WMF image — built as raw ZIP, no openpyxl image API needed)
# ---------------------------------------------------------------------------
def make_xlsx_with_wmf():
    """Build an XLSX that embeds a WMF image using raw ZIP construction.

    openpyxl silently drops WMF images on read, so we bypass the API and
    write the necessary OOXML parts by hand.  The WMF bytes are fake (null
    bytes) — we only detect by file extension, not content.
    """
    import zipfile
    import io

    # Minimal shared-strings and styles XMLs required by openpyxl to open the file.
    content_types = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Default Extension="wmf" ContentType="image/x-wmf"/>
  <Override PartName="/xl/workbook.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>
  <Override PartName="/xl/worksheets/sheet1.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>
  <Override PartName="/xl/sharedStrings.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sharedStrings+xml"/>
  <Override PartName="/xl/styles.xml" ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.styles+xml"/>
  <Override PartName="/xl/drawings/drawing1.xml" ContentType="application/vnd.openxmlformats-officedocument.drawing+xml"/>
</Types>"""

    root_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="xl/workbook.xml"/>
</Relationships>"""

    workbook_xml = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
          xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheets>
    <sheet name="XLSX-WMF-SHEET-v5w6x7y8" sheetId="1" r:id="rId1"/>
  </sheets>
</workbook>"""

    workbook_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" Target="worksheets/sheet1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/sharedStrings" Target="sharedStrings.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/styles" Target="styles.xml"/>
</Relationships>"""

    shared_strings = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<sst xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" count="2" uniqueCount="2">
  <si><t>XLSX-WMF-ROW-BEFORE</t></si>
  <si><t>XLSX-WMF-ROW-AFTER</t></si>
</sst>"""

    styles = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<styleSheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">
  <fonts count="1"><font><sz val="11"/><name val="Calibri"/></font></fonts>
  <fills count="2">
    <fill><patternFill patternType="none"/></fill>
    <fill><patternFill patternType="gray125"/></fill>
  </fills>
  <borders count="1"><border><left/><right/><top/><bottom/><diagonal/></border></borders>
  <cellStyleXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0"/></cellStyleXfs>
  <cellXfs count="1"><xf numFmtId="0" fontId="0" fillId="0" borderId="0" xfId="0"/></cellXfs>
</styleSheet>"""

    sheet1_xml = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main"
           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <sheetData>
    <row r="1"><c r="A1" t="s"><v>0</v></c></row>
    <row r="2"><c r="A2" t="s"><v>1</v></c></row>
  </sheetData>
  <drawing r:id="rId1"/>
</worksheet>"""

    sheet1_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/drawing" Target="../drawings/drawing1.xml"/>
</Relationships>"""

    # Drawing XML: one TwoCellAnchor referencing rId1 (our WMF)
    drawing_xml = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<xdr:wsDr xmlns:xdr="http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing"
           xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"
           xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">
  <xdr:twoCellAnchor>
    <xdr:from><xdr:col>0</xdr:col><xdr:colOff>0</xdr:colOff><xdr:row>0</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:from>
    <xdr:to><xdr:col>1</xdr:col><xdr:colOff>0</xdr:colOff><xdr:row>1</xdr:row><xdr:rowOff>0</xdr:rowOff></xdr:to>
    <xdr:pic>
      <xdr:nvPicPr>
        <xdr:cNvPr id="2" name="Picture 1"/>
        <xdr:cNvPicPr/>
      </xdr:nvPicPr>
      <xdr:blipFill>
        <a:blip r:embed="rId1"/>
        <a:stretch><a:fillRect/></a:stretch>
      </xdr:blipFill>
      <xdr:spPr>
        <a:xfrm><a:off x="0" y="0"/><a:ext cx="1000000" cy="1000000"/></a:xfrm>
        <a:prstGeom prst="rect"><a:avLst/></a:prstGeom>
      </xdr:spPr>
    </xdr:pic>
    <xdr:clientData/>
  </xdr:twoCellAnchor>
</xdr:wsDr>"""

    drawing_rels = b"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/image" Target="../media/image1.wmf"/>
</Relationships>"""

    # Fake WMF bytes (just a null placeholder — we only check the extension)
    wmf_bytes = b"\x00" * 8

    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", root_rels)
        zf.writestr("xl/workbook.xml", workbook_xml)
        zf.writestr("xl/_rels/workbook.xml.rels", workbook_rels)
        zf.writestr("xl/sharedStrings.xml", shared_strings)
        zf.writestr("xl/styles.xml", styles)
        zf.writestr("xl/worksheets/sheet1.xml", sheet1_xml)
        zf.writestr("xl/worksheets/_rels/sheet1.xml.rels", sheet1_rels)
        zf.writestr("xl/drawings/drawing1.xml", drawing_xml)
        zf.writestr("xl/drawings/_rels/drawing1.xml.rels", drawing_rels)
        zf.writestr("xl/media/image1.wmf", wmf_bytes)

    out_path = os.path.join(OUT, "test_xlsx_with_wmf.xlsx")
    with open(out_path, "wb") as f:
        f.write(buf.getvalue())
    print("Created test_xlsx_with_wmf.xlsx")


# ---------------------------------------------------------------------------
# XLSX (with two embedded images on the same sheet)
# ---------------------------------------------------------------------------
def make_xlsx_with_multiple_images():
    import struct
    import zlib
    import openpyxl
    from io import BytesIO
    from openpyxl.drawing.image import Image as XlImage

    def _minimal_png() -> bytes:
        def chunk(tag: bytes, data: bytes) -> bytes:
            crc = zlib.crc32(tag + data) & 0xFFFFFFFF
            return struct.pack(">I", len(data)) + tag + data + struct.pack(">I", crc)
        ihdr = chunk(b"IHDR", struct.pack(">IIBBBBB", 1, 1, 8, 2, 0, 0, 0))
        idat = chunk(b"IDAT", zlib.compress(b"\x00\xff\xff\xff"))
        iend = chunk(b"IEND", b"")
        return b"\x89PNG\r\n\x1a\n" + ihdr + idat + iend

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "XLSX-MULTI-SHEET-a1b2c3d4"

    ws["A1"] = "XLSX-MULTI-ROW-A"
    ws["A2"] = "XLSX-MULTI-ROW-B"
    ws["A3"] = "XLSX-MULTI-ROW-C"
    ws["A4"] = "XLSX-MULTI-ROW-D"

    # Image 1 anchored at row 1 (0-based) — splits between ROW-A and ROW-B
    img1 = XlImage(BytesIO(_minimal_png()))
    ws.add_image(img1, "A2")

    # Image 2 anchored at row 3 (0-based) — splits between ROW-C and ROW-D
    img2 = XlImage(BytesIO(_minimal_png()))
    ws.add_image(img2, "A4")

    wb.save(os.path.join(OUT, "test_xlsx_with_multiple_images.xlsx"))
    print("Created test_xlsx_with_multiple_images.xlsx")


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content

    slide1 = prs.slides.add_slide(layout)
    slide1.shapes.title.text = "PPTX-TITLE-d0e1f2g3"
    slide1.placeholders[1].text = "PPTX-BODY-h4i5j6k7"

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    txBox = slide2.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
    txBox.text_frame.text = "PPTX-SLIDE2-l8m9n0o1"

    prs.save(os.path.join(OUT, "test.pptx"))
    print("Created test.pptx")


# ---------------------------------------------------------------------------
# PDF
# ---------------------------------------------------------------------------
def make_pdf():
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=16)
    pdf.cell(0, 10, "PDF-HEADING-p2q3r4s5", ln=True)
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "PDF-PARA-t6u7v8w9", ln=True)
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(0, 10, "PDF-PAGE2-x0y1z2a3", ln=True)

    pdf.output(os.path.join(OUT, "test.pdf"))
    print("Created test.pdf")


if __name__ == "__main__":
    os.makedirs(OUT, exist_ok=True)
    make_docx()
    make_xlsx()
    make_xlsx_with_images()
    make_xlsx_with_wmf()
    make_xlsx_with_multiple_images()
    make_pptx()
    make_pdf()
    print("All test files created.")
